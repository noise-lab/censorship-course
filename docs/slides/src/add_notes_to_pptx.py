#!/usr/bin/env python3
"""
Script to add speaker notes to PowerPoint slides
"""

try:
    from pptx import Presentation
    import sys

    pptx_file = "11-PrivacyLaw.pptx"

    # Speaker notes for each slide (1-indexed)
    notes = {
        1: "",  # Title slide
        2: """Context: This is the 1890 Harvard Law Review article "The Right to Privacy"
What to say: "This article was written in response to new photography technology and gossip journalism. Warren was upset about newspapers covering his family's social events. Sound familiar? Privacy law always responds to new technology."
URL: https://www.cs.cornell.edu/~shmat/courses/cs5436/warren-brandeis.pdf""",

        3: """Context: 1970s concerns about mainframe computers
What to say: "These concerns from the 1970s sound quaint now, but they led to the Fair Information Practice Principles. Notice how every concern is 100x worse today with cloud computing and smartphones." """,

        4: """Context: 1973 HEW report
What to say: "These principles became the foundation for most privacy law worldwide. Notice the focus on transparency and individual control - we'll come back to whether this model still works." """,

        5: "",  # Empty
        6: "",  # Empty

        7: """What to say: "This is why every website has a privacy policy. Problem: they're unreadable. Instagram's is 10,000+ words. The notice-and-consent model is fundamentally broken."
URL to demo: https://tosdr.org/ - "Terms of Service; Didn't Read" rates privacy policies""",

        8: """What to say: "Primary use: obvious. Secondary uses: contentious. When is it OK to use data for purposes beyond what it was collected for?"
Example: Target figured out a teen was pregnant before her father knew, based on purchasing patterns. Appropriate use of data?""",

        9: """What to say: "Classic secondary use debate. Collected for employment, used for law enforcement. Most people think this is reasonable. But where's the line?" """,

        10: """What to say: "GDPR requires opt-in for most things. US typically allows opt-out. Huge difference in practice - opt-out means most people never opt out."
Demo: Show annoying cookie consent banners on EU sites vs. US sites with no banners""",

        11: """What to say: "GDPR's 'Right to Access' is powerful - you can request all data a company has on you. Try it with Google or Facebook - you'll be shocked."
URL: https://support.google.com/accounts/answer/3024190 - How to download your Google data""",

        12: """What to say: "Security enables privacy. No point in privacy laws if data gets breached. This led to breach notification laws." """,

        13: """What to say: "Modern additions: minimization (collect only what you need), breach notification, accountability. These evolved as weaknesses in original FIPPs became clear." """,

        14: """What to say: "Are FIPPs outdated? Example: differential privacy deliberately adds noise to data. That violates 'data must be correct' but enhances privacy. The paradigm may need updating for modern ML/AI." """,

        15: """What to say: "Old model: data is static, stored in databases. New model: data is queried, inferred from, probabilistically analyzed. FIPPs weren't designed for this." """,

        16: """What to say: "PII concept is also outdated. Remove name, SSN, address = safe? Wrong. Netflix 'anonymized' data was re-identified. AOL search data was re-identified. PII is not a bright line."
Famous case: Netflix Prize dataset re-identification (2006)""",

        17: """What to say: "FIPPs assume data sits in databases and gets retrieved. But modern ML does inference, prediction, pattern matching. Different threat model." """,

        18: "",  # Image slide
        19: "",  # Image slide

        20: """What to say: "Modern laws like California's CCPA try to fix the PII problem with 'de-identification' standards. But it's hard to define 'reasonably linked.'" """,

        21: """What to say: "US privacy law is a patchwork based on: (1) who has the data, (2) what harm might occur, (3) what transaction is happening. No comprehensive federal law." """,

        22: """What to say:
- Privacy Act (1974): Post-Watergate, only applies to federal agencies
- FCRA (1970): Credit bureaus only
- FERPA (1974): Schools collected data on kids' home lives - Congress freaked out
- Video Privacy Protection Act (1988): Reporter published Judge Bork's video rental history during Supreme Court nomination. Congress passed law within months.
- Driver's Privacy Protection Act (1994): States were selling DMV data""",

        23: """What to say:
- HIPAA (1996): Covers healthcare providers/insurers but NOT fitness apps, 23andMe, etc.
- GLBA (1999): Banks can share data among affiliates
- COPPA (1998): Kids under 13 need parental consent
- GINA (2008): Can't discriminate based on genetic information""",

        24: """What to say: "Huge gap: fitness apps, period trackers, mental health apps are NOT covered by HIPAA. They can sell your health data."
Recent case: BetterHelp (mental health app) fined $7.8M for sharing data with Facebook
URL: https://www.ftc.gov/news-events/news/press-releases/2023/03/ftc-ban-betterhelp-revealing-consumers-data-including-sensitive-mental-health-information-facebook""",

        25: """What to say: "Before 2003, companies could suffer massive breaches and never tell anyone. California changed everything." """,

        26: """What to say: "Now all 50 states have breach notification laws. Made security a board-level issue. If you get breached, it's front page news." """,

        27: """What to say: "Varies by state - this is the patchwork problem. Some states require notification for any breach, others only if 'harm likely.'" """,

        28: """What to say: "Companies just apply the strictest state law (usually California) to everyone. Easier than 50 different compliance regimes." """,

        29: """What to say: "147 million Americans' data breached. Equifax's response was a disaster: took 6 weeks to disclose, set up fake-looking website, tried to profit from credit monitoring."
URL: https://www.ftc.gov/equifax-data-breach
Result: $700M settlement, CEO resigned""",

        30: """What to say: "FTC is the de facto privacy regulator in the US, but they don't have a specific privacy statute. They use Section 5..." """,

        31: """What to say: "Section 5 bans 'unfair or deceptive' practices. FTC stretches this to cover privacy. Not ideal - built for consumer protection, not privacy rights."
Deceptive: You promised X, did Y
Unfair: Causes substantial harm consumers can't avoid""",

        32: """What to say: "Uber claimed 'God View' tracking of users was limited. FTC found they were monitoring journalists, politicians, exes. Settled with consent decree."
URL: https://www.ftc.gov/news-events/news/press-releases/2017/08/uber-settles-ftc-allegations-it-made-deceptive-privacy-data-security-claims""",

        33: """Context: 2011 consent decree, then Cambridge Analytica 2018
What to say: "Facebook settled with FTC in 2011 for privacy violations. Consent decree required audits, compliance. Then Cambridge Analytica happened - 87 million users' data harvested. Did Facebook violate the decree?"
URL: https://www.ftc.gov/news-events/news/press-releases/2019/07/ftc-imposes-5-billion-penalty-sweeping-new-privacy-restrictions-facebook
Result: $5 billion fine (2019) - largest FTC penalty ever""",

        34: "",  # Image about Facebook

        35: """What to say: "2019 consent decree is much stricter. Two-factor authentication phone numbers were being used for ads - now prohibited. Facial recognition requires explicit consent." """,

        36: """What to say: "140+ countries have privacy laws. EU's GDPR is the most influential - it has extraterritorial reach, so it affects US companies too." """,

        37: """What to say: "GDPR applies to ANY company offering services to EU residents. US companies must comply. This is regulatory imperialism - EU law governing US companies."
Example: If you have a website that EU citizens can access, GDPR may apply to you.""",

        38: """What to say: "'Right to be forgotten' is controversial. Google must remove search results in EU if requested (with exceptions). Doesn't apply in US - First Amendment conflict."
URL: https://transparencyreport.google.com/eu-privacy/overview - Google's transparency report on GDPR requests""",

        39: """What to say: "Cross-border transfers are complicated. EU says US doesn't have 'adequate' privacy protection. Privacy Shield was invalidated in 2020 (Schrems II). New framework in 2023, but legal challenges continue."
URL: https://www.dataprivacyframework.gov/ - New EU-US Data Privacy Framework""",

        40: """Context: GDPR enforcement examples
What to show: https://www.enforcementtracker.com/ - Live GDPR fines tracker
Examples: Meta €1.2B (2023), Amazon €746M (2021), Google €90M (2021)""",

        41: """What to say: "Debate: Does GDPR protect privacy or just make compliance so expensive that only big tech can afford it? Small startups struggle with compliance. Cookie consent is annoying UX. But enforcement is real."
Counterpoint: GDPR hasn't stopped big tech dominance""",

        42: """URLs to show:
- TikTok: https://www.nytimes.com/2024/01/29/technology/tiktok-data-privacy-lawsuit.html
- Cars selling data: https://foundation.mozilla.org/en/privacynotincluded/categories/cars/
- GoodRx: https://www.ftc.gov/news-events/news/press-releases/2023/02/ftc-enforcement-action-bar-goodrx-sharing-consumers-sensitive-health-info-advertising
- Twitter 2FA: https://www.ftc.gov/news-events/news/press-releases/2022/05/ftc-charges-twitter-deceptively-using-account-security-data-sell-targeted-ads"""
    }

    prs = Presentation(pptx_file)

    print(f"Adding speaker notes to {len(notes)} slides...\n")

    for slide_num, note_text in notes.items():
        if note_text:  # Only add if there's actual content
            slide_idx = slide_num - 1  # Convert to 0-indexed
            if slide_idx < len(prs.slides):
                slide = prs.slides[slide_idx]

                # Get or create notes slide
                if not slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                else:
                    notes_slide = slide.notes_slide

                # Set the notes text
                text_frame = notes_slide.notes_text_frame
                text_frame.text = note_text

                print(f"✓ Added notes to slide {slide_num}")

    # Save the presentation
    output_file = pptx_file
    prs.save(output_file)
    print(f"\n✅ Saved presentation to: {output_file}")

except ImportError:
    print("python-pptx not installed. Installing...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    print("\nPlease run the script again.")
    sys.exit(1)
except Exception as e:
    print(f"Error: {e}")
    import traceback
    traceback.print_exc()
    sys.exit(1)
