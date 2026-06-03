#!/usr/bin/env python3
"""
Extract per-slide text, titles, and speaker notes from a .pptx into a deck
folder's _source-extract.md, and dump every embedded image to images/.

Usage:
    python read_pptx.py <deck.pptx> <out-dir>

<out-dir> is created if needed; writes <out-dir>/_source-extract.md and
<out-dir>/images/slideNNN_imgMMM.<ext>. These are extraction inputs only
(gitignored); the authored slides.qmd references just the useful images.
"""
import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Emu  # noqa: F401  (imported for completeness)
except ImportError:
    sys.exit("python-pptx not installed. Run: pip install python-pptx")


def extract(pptx_file, out_dir):
    prs = Presentation(pptx_file)
    img_dir = os.path.join(out_dir, "images")
    os.makedirs(img_dir, exist_ok=True)

    lines = []
    lines.append(f"# Source extract: {os.path.basename(pptx_file)}")
    lines.append(f"\nTotal slides: {len(prs.slides)}\n")

    n_imgs = 0
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n## SLIDE {i}")
        if slide.shapes.title and slide.shapes.title.text.strip():
            lines.append(f"\n**Title:** {slide.shapes.title.text.strip()}")

        body = []
        for shape in slide.shapes:
            # text
            if shape.has_text_frame and shape.text.strip():
                if slide.shapes.title and shape == slide.shapes.title:
                    continue
                body.append(shape.text.strip())
            # images
            if shape.shape_type == 13 or getattr(shape, "image", None) is not None:
                try:
                    image = shape.image
                    n_imgs += 1
                    ext = image.ext
                    fname = f"slide{i:03d}_img{n_imgs:03d}.{ext}"
                    with open(os.path.join(img_dir, fname), "wb") as fh:
                        fh.write(image.blob)
                    lines.append(f"\n_[image: images/{fname}]_")
                except Exception:
                    pass

        if body:
            lines.append("\n**Content:**")
            for b in body:
                for ln in b.splitlines():
                    if ln.strip():
                        lines.append(f"- {ln.strip()}")

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"\n**Speaker notes:**\n\n> {notes}")

    out_md = os.path.join(out_dir, "_source-extract.md")
    with open(out_md, "w") as fh:
        fh.write("\n".join(lines) + "\n")

    print(f"{os.path.basename(pptx_file)}: {len(prs.slides)} slides, "
          f"{n_imgs} images -> {out_dir}")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        sys.exit(__doc__)
    extract(sys.argv[1], sys.argv[2])
